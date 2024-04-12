from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd

# 创建一个示例 DataFrame
data = {
    'ID': [101, 102, 103],
    'Name': ['Alice (moved)', 'Bob (new)', 'Charlie']
}
df = pd.DataFrame(data)

# 创建一个 PowerPoint 对象
prs = Presentation()

# 在 PowerPoint 中创建一个新的幻灯片
slide_layout = prs.slide_layouts[5]  # 选择幻灯片布局，这里选择了标题和内容
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title = slide.shapes.title
title.text = "DataFrame to PowerPoint Example"

# 添加一个空白幻灯片
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 添加表格
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(3)
table = slide.shapes.add_table(df.shape[0] + 1, df.shape[1], left, top, width, height).table

# 设置表格标题行的背景色和字体样式，并填充为 DataFrame 的列名
for idx, column_name in enumerate(df.columns):
    cell = table.cell(0, idx)
    cell.text = column_name
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 设置字体颜色为白色
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 128, 0)  # 设置背景色为深绿色

# 添加数据
for i in range(df.shape[0]):
    for j in range(df.shape[1]):
        cell = table.cell(i+1, j)
        cell.text = str(df.iloc[i, j])
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # 设置字体颜色为黑色

# 设置内容行的背景色为浅绿和浅蓝交替
for i in range(1, df.shape[0] + 1):
    if i % 2 == 0:
        for cell in table.rows[i].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(173, 216, 230)  # 设置背景色为浅蓝色
    else:
        for cell in table.rows[i].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(144, 238, 144)  # 设置背景色为浅绿色

# 保存 PowerPoint 文件
prs.save('dataframe_to_ppt_example.pptx')
