import pandas as pd
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
#from pptx.enum.shapes import MSO_SHAPE, MSO_ANCHOR
import argparse

parser = argparse.ArgumentParser(description='python ConvertToPPTX.py --f MMYY_output_SrcPPT_YoY.xlsx --s Sheet1')
parser.add_argument('--f', help='MMYY_output_SrcPPT_YoY.xlsx')
parser.add_argument('--s', help='Sheet1')
args = parser.parse_args()

fileName = args.f
sheetName = args.s
# 创建一个示例 DataFrame
df = pd.read_excel(fileName, sheet_name=sheetName)
#df = pd.read_excel("MMYY_output_SrcPPT_YoY.xlsx", sheet_name="Sheet1")

# 创建一个 PowerPoint 对象
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 在 PowerPoint 中创建一个新的幻灯片
slide_layout = prs.slide_layouts[5]  # 选择幻灯片布局，这里选择了标题和内容
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title = slide.shapes.title
title.text = "DataFrame to PowerPoint Example"

# 添加一个空白幻灯片
blank_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(blank_slide_layout)

# 添加表格
left = Inches(0.5)
top = Inches(1.5)
width = Inches(15)
height = Inches(3)
table = slide.shapes.add_table(df.shape[0] + 1, df.shape[1], left, top, width, height).table

Font_Title_Size = Pt(13)
Font_Title_Bold = True
RGB_Title_Text = RGBColor(255, 255, 255)  # 设置字体颜色为白色
RGB_Title_Text_Mark_GOLD = RGBColor(255, 255, 0)  # 设置字体颜色为金黄色
RGB_Title_Text_Mark_RED = RGBColor(255, 0, 0)  # 设置字体颜色为红色
RGB_Title_Cell = RGBColor(29, 154, 120)  # 设置背景色为深绿色

Font_Content_Size = Pt(11)
RGB_Content_Text = RGBColor(0, 0, 0)  # 设置字体颜色为黑色
RGB_Content_Text_Mark_Negative = RGBColor(255, 0, 0)  # 设置字体颜色为红色
RGB_Content_Text_Mark_Positive = RGBColor(0, 128, 0)  # 设置字体颜色为深绿色
RGB_Content_Cell_LightGreen = RGBColor(204, 222, 214)  # 设置背景色为浅绿色
RGB_Content_Cell_LightBlue = RGBColor(231, 239, 236)  # 设置背景色为浅蓝色

title_MarkGold_Strs = ['MoM', 'YoY', '%', '@ Jan', '@ Feb', '@ Mar', '@ Apr', '@ May', '@ Jun', '@ Jul', '@ Aug', '@ Sep', '@ Oct', '@ Nov', '@ Dec']
title_MarkRed_Strs = ['10']
content_Red_Strs = ['moved', 'decreased']
content_Green_Strs = ['new', 'increased', 'returned']

# 设置指定单元格的某些字符的颜色, 这些字符不能有包含或重复关系
def ColorCellPartTexts(cell, part_texts1, color_pain1, part_texts2, color_pain2):
    print("!!!!!!!!!!!!! ColorCellPartTexts")
    
    if len(cell.text) == 0:
        return
    
    if len(part_texts1) == 0:
        return
    
    cell_txt = cell.text
    len_cell_text = len(cell_txt)
    dict_idxPairToColor = {}
    
    if "Corp-83624040-SBR_SoftBank_Affiliate" in cell_txt :
        print("Find: ", "Corp-83624040-SBR_SoftBank_Affiliate")
    
    # 保存字体样式
    cur_alignment = cell.text_frame.paragraphs[0].alignment
    font_Bold = cell.text_frame.paragraphs[0].font.bold
    font_Size = cell.text_frame.paragraphs[0].font.size
    font_color = cell.text_frame.paragraphs[0].font.color.rgb
    
    for part_text in part_texts1:
        start_pos = 0
        while True:
            idx = cell_txt.find(part_text, start_pos)
            if idx == -1:
                break
            key = (idx, idx+len(part_text))
            dict_idxPairToColor[key] = color_pain1
            start_pos = idx+len(part_text)
            print(f"^^^^^ Find: {part_text} inStr: {cell_txt[key[0]: key[1]]} idx: {idx} start_pos: {start_pos} len(cell_txt): {len(cell_txt)} color: {color_pain1}")
        
    if part_texts2 and len(part_texts2) > 0:
        for part_text in part_texts2:
            start_pos = 0
            while True:
                idx = cell_txt.find(part_text, start_pos)
                if idx == -1:
                    break
                key = (idx, idx+len(part_text))
                dict_idxPairToColor[key] = color_pain2
                start_pos = idx+len(part_text)
                print(f"^^^^^ Find: {part_text} inStr: {cell_txt[key[0]: key[1]]} idx: {idx} start_pos: {start_pos} len(cell_txt): {len(cell_txt)} color: {color_pain2}")
            
    if len(dict_idxPairToColor) == 0 :
        return
    
    sorted_dict = dict(sorted(dict_idxPairToColor.items(), key=lambda x: x[0][0]))
    dict_idxPairToColor = sorted_dict
    
    lst_parts_color = []
    lst_dict_KeyPairs = list(dict_idxPairToColor.keys())
    len_dict_KeyPairs = len(lst_dict_KeyPairs)
    cur_part_id = lst_dict_KeyPairs[0][0]
    
    # 前面有其他字符
    if cur_part_id > 0:
        part_text = cell_txt[0:cur_part_id]
        lst_parts_color.append((part_text, font_color))
        
    # 中间有需要标记的字符
    pairis_Index = 0
    while pairis_Index < len_dict_KeyPairs:
        this_pair = lst_dict_KeyPairs[pairis_Index]
        # 正常需要标记的字符
        part_text = cell_txt[this_pair[0]:this_pair[1]]
        lst_parts_color.append((part_text, dict_idxPairToColor[this_pair]))
        # 需要标记的字符之间的其他字符
        if pairis_Index < len_dict_KeyPairs - 1:
            next_pairs = lst_dict_KeyPairs[pairis_Index + 1]
            # 两个标记字符之间有其他字符
            if this_pair[1] < next_pairs[0] :
                part_text = cell_txt[this_pair[1]:next_pairs[0]]
                lst_parts_color.append((part_text, font_color))
        pairis_Index = pairis_Index + 1
        
    # 后面有其他字符
    lastMarkedPairEndValue = lst_dict_KeyPairs[len_dict_KeyPairs - 1][1]
    if lastMarkedPairEndValue < len_cell_text:
        part_text = cell_txt[lastMarkedPairEndValue:]
        lst_parts_color.append((part_text, font_color))
    
    print(f"====== cell_txt: {cell_txt}")
    print(f"====== lst_parts_color: {lst_parts_color}")
    
    cell.text_frame.clear()
    p = cell.text_frame.paragraphs[0]
    p.alignment = cur_alignment
    for part_text, color in lst_parts_color:
        run = p.add_run()
        run.text = part_text
        run.font.bold = font_Bold
        run.font.size = font_Size
        run.font.color.rgb = color
    print(cell.text, " ", len(cell.text_frame.paragraphs))

def Merge1stColumn(table):
    total_rows = len(table.rows)
    
    row_idx = 1
    while row_idx < total_rows - 1:
        toMerge = False    
        cell = table.cell(row_idx, 0)  # 第一列
        row2_idx = row_idx + 1
        while row2_idx < total_rows:
            cell2 = table.cell(row2_idx, 0)
            if cell.text == cell2.text:
                #print("Merge row_idx: ", row_idx, " row2_idx: ", row2_idx)
                toMerge = True
                row2_idx = row2_idx + 1
            else:
                break
        #print("toMerge: ", toMerge, " row_idx: ", row_idx, " ToMerge: ", row2_idx - 1)
    
        if toMerge:
            #print("Merge row_idx: ", row_idx, " ToMerge: ", row2_idx - 1)
            cell.merge(table.cell(row2_idx - 1, 0))
            cell.text = cell.text_frame.paragraphs[0].text
            cell.text_frame.paragraphs[0].font.size = Font_Content_Size
            
        row_idx = row2_idx

# 设置表格标题行的背景色和字体样式
#for cell in table.rows[0].cells:
for idx, column_name in enumerate(df.columns):
    cell = table.cell(0, idx)
    cell.text = column_name
    cell.text_frame.paragraphs[0].font.bold = Font_Title_Bold
    cell.text_frame.paragraphs[0].font.size = Font_Title_Size
    cell.text_frame.paragraphs[0].font.color.rgb = RGB_Title_Text  # 设置字体颜色为白色
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGB_Title_Cell  # 设置背景色为深绿色
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER # 将文本水平和垂直对齐方式设置为居中
    ColorCellPartTexts(cell, title_MarkGold_Strs, RGB_Title_Text_Mark_GOLD, title_MarkRed_Strs, RGB_Title_Text_Mark_RED)

table.columns[0].width = Inches(15 / 8 * 0.8)
table.columns[1].width = Inches(15 / 8 * 0.9)
table.columns[2].width = Inches(15 / 8 * 0.5)
table.columns[3].width = Inches(15 / 8 * 0.4)
table.columns[4].width = Inches(15 / 8 * 0.3)
table.columns[5].width = Inches(15 - (15 / 8 * 3.7))
table.columns[6].width = Inches(15 / 8 * 0.4)
table.columns[7].width = Inches(15 / 8 * 0.4)

# 添加数据
for i in range(df.shape[0]):
    for j in range(df.shape[1]):
        cell = table.cell(i+1, j)
        cell.text = str(df.iloc[i, j])
        cell.text_frame.paragraphs[0].font.size = Font_Content_Size
        cell.text_frame.paragraphs[0].font.color.rgb = RGB_Content_Text  # 设置字体颜色为黑色
        ColorCellPartTexts(cell, content_Green_Strs, RGB_Content_Text_Mark_Positive, content_Red_Strs, RGB_Content_Text_Mark_Negative)
        # 如果 "Name" 列中包含 "moved" 字符串，将其标记为红色
        #if isinstance(df.iloc[i, j], str) and "(moved)" in df.iloc[i, j]:
        #    idx = df.iloc[i, j].find("(moved)") + 1
        #    part1 = df.iloc[i, j][:idx]
        #    part2 = df.iloc[i, j][idx:idx+len("(moved)")-2]
        #    part3 = df.iloc[i, j][idx+len("(moved)")-2:]
        #    #print(part1)
        #    #print(part2)
        #    #print(part3)
        #    cell.text_frame.clear()  # 清空单元格中原有的内容
        #    p = cell.text_frame.add_paragraph()
        #    p.text = part1
        #    run = p.add_run()
        #    run.text = part2
        #    font = run.font
        #    font.color.rgb = RGBColor(255, 0, 0)  # 设置字体颜色为红色
        #    #p = cell.text_frame.add_paragraph()
        #    #p.text = part3
        #    run2 = p.add_run()
        #    run2.text = part3
        #    cell.text_frame.paragraphs[0].font.size = Pt(10)
        #    for para in cell.text_frame.paragraphs:
        #        para.font.size = Pt(10)
        ## 如果 "Name" 列中包含 "new" 字符串，将其标记为深绿色加粗
        #elif isinstance(df.iloc[i, j], str) and "(new)" in df.iloc[i, j]:
        #    idx = df.iloc[i, j].find("(new)")
        #    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 128, 0)  # 设置字体颜色为深绿色
        #    cell.text_frame.paragraphs[0].runs[0].font.bold = True
        #    for para in cell.text_frame.paragraphs:
        #        para.font.size = Pt(10)

    
    #next_cell = table.cell(row_idx + 1, 0)  # 下一行的第一列
    ## 检查当前单元格和下一行单元格的文本内容是否相同
    #if cell.text == next_cell.text:
    #    # 如果相同，合并当前单元格和下一行单元格
    #    cell.merge(next_cell)
    #    cell.text = cell.text_frame.paragraphs[0].text
    #    cell.text_frame.paragraphs[0].font.size = Pt(10)




# 设置内容行的背景色为浅绿和浅蓝交替
for i in range(1, df.shape[0] + 1):
    if i % 2 == 0:
        for cell in table.rows[i].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGB_Content_Cell_LightBlue  # 设置背景色为浅蓝色
    else:
        for cell in table.rows[i].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGB_Content_Cell_LightGreen  # 设置背景色为浅绿色

Merge1stColumn(table)
# 保存 PowerPoint 文件
#prs.save('MMYY_output_SrcPPT_YoY.pptx')
prs.save(fileName.replace(".xlsx", ".pptx"))
